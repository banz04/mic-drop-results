import contextlib
from ctypes import windll
from io import BytesIO
import itertools
from multiprocessing import Pool, freeze_support
import os
import re
from signal import signal, SIGINT, SIG_IGN
from subprocess import check_call, run, DEVNULL
import sys
import time
import webbrowser

from colorama import init, Fore, Style
import cursor
import cv2
import numpy as np
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_SHAPE  # type: ignore
from pptx.slide import Slide
from pptx.util import Inches
from pywintypes import com_error
import requests
import win32com.client

from client import ProgramStatus, fetch_latest_version
from client import download_avatar, get_avatar_path
from config import Config
from constants import *
from errors import Error, ErrorType, print_exception_hook
from exceptions import *
from utils import is_number, as_type, hex_to_rgb, parse_version
from utils import abs_path
from utils import inp, console_style, ProgressBar
from vba.macros import module1_bas


def replace_text(slide: Slide, df, i, avatar_mode) -> Slide:
    """Replaces and formats text."""
    cols = df.columns.tolist() + ['p']

    for shape in slide.shapes:  # type: ignore
        if not shape.has_text_frame or '{' not in shape.text:
            continue

        text_frame = shape.text_frame

        for run in itertools.chain.from_iterable(
            [p.runs for p in text_frame.paragraphs]):

            for search_str in (set(re.findall(r'(?<={)(.*?)(?=})', run.text))
                               .intersection(cols)):

                # Avatars
                if search_str == 'p':
                    # Test cases
                    if '__uid__' not in cols or not avatar_mode:
                        run.text = ''
                        continue

                    if pd.isnull(df['__uid__'].iloc[i]):
                        run.text = ''
                        continue

                    # Extract effect index and remove {p}
                    effect = run.text.strip()[3:]
                    run.text = ''

                    uid: str = df['__uid__'].iloc[i]

                    avatar_path = get_avatar_path(AVATAR_DIR, uid)
                    avatarfx_path = get_avatar_path(AVATAR_DIR, uid,
                                                    effect=effect)

                    if not os.path.isfile(avatar_path):
                        continue

                    if is_number(effect):
                        img = cv2.imread(avatar_path)
                        match float(effect):  # TODO: Add more effects
                            case 1:
                                img = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

                        cv2.imwrite(avatarfx_path, img)

                    new_shape = slide.shapes.add_picture(  # type: ignore
                        avatarfx_path, shape.left, shape.top,
                        shape.width, shape.height
                    )

                    new_shape.auto_shape_type = MSO_SHAPE.OVAL
                    old = shape._element
                    new = new_shape._element
                    old.addnext(new)
                    old.getparent().remove(old)
                    continue

                # Actual text
                repl = str(df[search_str].iloc[i])
                repl = repl if repl != 'nan' else ''  # Replace nan with empty

                run_text = run.text

                if search_str.startswith(cfg.trigger_word):
                    run.text = repl
                else:
                    run.text = run.text.replace('{' + search_str + '}', repl)

                # Replace image links
                pattern = r'\<\<(.*?)\>\>'  # Look for <<image_links.{ext}>>
                img_link = re.findall(pattern, run.text)

                if len(img_link) > 0:
                    try:
                        img = BytesIO(requests.get(img_link[0]).content)
                        pil = Image.open(img)

                        im_width = shape.height / pil.height * pil.width
                        new_shape = slide.shapes.add_picture(  # type: ignore
                            img, shape.left + (shape.width - im_width)/2,
                            shape.top, im_width, shape.height
                        )

                        old = shape._element.addnext(new_shape._element)

                        run.text = re.sub(pattern, '', run.text)
                        text_frame.margin_left = Inches(5.2)
                    except Exception:
                        Error(
                            'Could not load the following image from '
                           f'slide {i + 1}, sheet {df["sheet"].iloc[0]}.',
                           f'{img_link[0]}',
                            'Please check your internet connection and verify '
                            'that the link directs to an image file, which '
                            'usually ends in an image extension like .png.',
                            err_type=ErrorType.WARNING).throw()

                # Color formatting
                if not search_str.startswith(cfg.trigger_word):
                    continue

                # Check RGB
                if (run.font.color.type == MSO_COLOR_TYPE.RGB and  # type: ignore
                    run.font.color.rgb not in [
                        RGBColor(0, 0, 0), RGBColor(255, 255, 255)]):
                    continue

                for ind, val in enumerate(cfg.ranges[::-1]):
                    if is_number(repl) and float(repl) >= val:
                        if run_text.endswith('1'):
                            run.font.color.rgb = RGBColor(
                                *scheme_alt[::-1][ind])
                        else:
                            run.font.color.rgb = RGBColor(
                                *scheme[::-1][ind])
                        break
    return slide


def preview_df(df: pd.DataFrame, filter_series: pd.Series | None = None, *,
               n_cols: int, n_cols_ext: int = 5,
               highlight: bool = True,
               words_to_highlight: list[str | None] | None = None) -> str:

    if words_to_highlight is None:
        words_to_highlight = []

    df.index += 2  # To reflect row numbers as displayed in Excel

    # Only show the first few columns for preview
    df = df.iloc[:, : min(n_cols + n_cols_ext, len(df.columns))]

    if filter_series:
        filter_series.index += 2
        df = df[filter_series]

    # Replace values_to_highlight with ???values_to_highlight???
    # The brackets are weird Unicode characters that no one would ever use.
    for word in words_to_highlight:
        if word is None:
            df.iloc[:, :n_cols] = df.iloc[:, :n_cols].fillna('???NaN???')
        else:
            df.iloc[:, :n_cols] = df.iloc[:, :n_cols].replace(
                word, f'???{word}???')


    preview = repr(df.head(8))

    # Highlight ???values_to_highlight???
    preview = preview.replace('???', Fore.RED + '  ').replace('???', Fore.RESET)

    # Highlight column names
    for n in range(n_cols):
        col = df.columns.tolist()[n]
        preview = preview.replace(
            ' ' + col, f' {Fore.RED}{col}{Fore.RESET}', 1)


    # Add ... at the end of each line if preview is a snippet
    if n_cols < len(df.columns):
        preview = preview.replace('\n', '  ...\n') + '  ...'

    # Bold first row
    preview = f'{Style.BRIGHT}' + preview.replace('\n', Style.NORMAL + '\n', 1)


    if not highlight:
        preview = preview.replace(Fore.RED, '')
    return preview


if __name__ == '__main__':
    version_tag = '2.3.0'

# Section A: Fix console issues
    freeze_support()          # Multiprocessing freeze support
    signal(SIGINT, SIG_IGN)   # Handle KeyboardInterrupt

    # Disable QuickEdit and Insert mode
    kernel32 = windll.kernel32
    kernel32.SetConsoleMode(
        kernel32.GetStdHandle(-10), (0x4|0x80|0x20|0x2|0x10|0x1|0x00|0x100))

    sys.excepthook = print_exception_hook  # Avoid exiting program on exception
    init()                                 # Enable ANSI escape sequences
    cursor.hide()                          # Hide cursor


# Section B: Check for missing files
    if missing := [f for f in (
            'settings.ini',
            'data.xlsx',
            'template.pptm',
            'token.txt',
        ) if not os.path.exists(abs_path(f))]:
        Error(40).throw(APP_DIR, '\n'.join(missing))


# Section C: Load user settings
    cfg = Config(abs_path('settings.ini'))

    n_scols = len(cfg.sorting_orders)

    avatar_mode = cfg.avatar_mode

    scheme = list(map(hex_to_rgb, cfg.scheme))
    scheme_alt = list(map(hex_to_rgb, cfg.scheme_alt))


# Section D: Parse and test tokens
    with open(abs_path('token.txt')) as f:
        token_list = f.read().splitlines()
        token_list = [line.replace('"', '').strip()
                      for line in token_list if len(line) > 70]

    if not token_list and avatar_mode:
        Error(21).throw()


# Section E: Check for updates
    status = None
    with contextlib.suppress(requests.exceptions.ConnectionError, KeyError):
        if cfg.update_check:
            # Fetch the latest version and the summary of the update
            latest_tag, summary = fetch_latest_version()

            latest, current = parse_version(
                latest_tag, version_tag
            )

            if latest > current:
                status = ProgramStatus.UPDATE_AVAILABLE

                console_style(Fore.YELLOW, Style.BRIGHT)
                print(f'Update v{latest_tag}')

                console_style(Style.NORMAL)
                print(summary)

                print(LATEST_RELEASE_URL)
                print()
                console_style()

                webbrowser.open(LATEST_RELEASE_URL, new=2)

            elif latest < current:
                status = ProgramStatus.BETA
            else:
                status = ProgramStatus.UP_TO_DATE

    # Print a header containing information about the program

    # Normal:       Mic Drop Results (v3.10) [latest]
    #               https://github.com/banz04/mic-drop-results


    # With update:  Update v3.11
    #               A summary of the update will appear in this line.
    #               https://github.com/banz04/mic-drop-results/releases/latest/
    #
    #               Mic Drop Results (v3.10) [update available]

    console_style(Style.BRIGHT)
    status_msg = f' [{status.value}]' if status else ''
    print(f'Mic Drop Results (v{version_tag}){status_msg}')
    console_style()

    if status != ProgramStatus.UPDATE_AVAILABLE:
    # When an update is available, the download link is already shown above.
    # To avoid confusion, we only print one link at a time.
        print(REPO_URL)


# Section F: Read and process data.xlsx
    xls = pd.ExcelFile(abs_path('data.xlsx'))
    workbook: dict[int | str, pd.DataFrame] = pd.read_excel(
        xls, sheet_name=None)
    xls.close()

    sheet_names = [re.sub(
        r'[\\\/:"*?<>|]+', '',  # Forbidden file name characters
        str(name)).strip() for name in workbook]

    workbook = {name: list(workbook.values())[i]
                for i, name in enumerate(sheet_names)}


    # Extract tables that belong to the database
    database: dict[str, pd.DataFrame] = {}
    for sheet in sheet_names:
        if not sheet.startswith('_'):  # Database tables are signified with _
            continue

        table = workbook[sheet]
        if table.empty or table.shape < (1, 2):  # (1 row, 2 cols) min
            continue
        database[sheet] = table


    # Process non-database sheets
    groups: dict[str, pd.DataFrame] = {}
    for sheet in sheet_names:
        if sheet.startswith('_'):  # Exclude database tables
            continue

        df = workbook[sheet]
        if df.empty or df.shape < (1, n_scols):  # (1 row, n sorting cols) min
            continue

        # Get sorting columns
        scols = df.columns.tolist()[:n_scols]
        SORTING_COLUMNS = (
            f'{Style.BRIGHT}Sheet name:{Style.NORMAL}       {sheet}\n'
            f'{Style.BRIGHT}Sorting columns:{Style.NORMAL}  {", ".join(scols)}'
            f'\n\nPlease have a look at the following rows in data.xlsx '
            f'to find out what caused the problem.'
        )


        # Exclude sheets with non-numeric sorting columns
        if any(df.loc[:, scol].dtype.kind not in 'biufc' for scol in scols):
            # Get list of non-numeric values
            str_vals = (df.loc[:, scols][~df.loc[:, scols].applymap(np.isreal)]
                .melt(value_name='__value__').dropna()['__value__'].tolist())

            Error(60).throw(
                SORTING_COLUMNS,

                preview_df(
                    df, ~df.loc[:, scols].applymap(np.isreal).all(1),
                    n_cols=n_scols, words_to_highlight=str_vals),

                err_type=ErrorType.WARNING)

            continue

        # Fill nan values within the sorting columns
        if df.loc[:, scols].isnull().values.any():
            print(df.loc[:, scols].isnull().any(axis=1))
            Error(61).throw(
                SORTING_COLUMNS,

                preview_df(
                    df, df.loc[:, scols].isnull().any(axis=1),
                    n_cols=n_scols, words_to_highlight=[None]),

                err_type=ErrorType.WARNING)

            df.loc[:, scols] = df.loc[:, scols].fillna(0)

        # Rank data
        df['__rank__'] = (
            pd.DataFrame(df.loc[:, scols]            # Select sorting columns
            * (np.array(cfg.sorting_orders)*2 - 1))  # Turn bool 0/1 into -1/1
            .apply(tuple, axis=1)  # type: ignore
            .rank(method='min', ascending=False)
            .astype(int))

        # Sort the slides
        df = df.sort_values(by='__rank__', ascending=True)

        # Remove .0 from whole numbers
        format_int = lambda x: str(int(x)) if x % 1 == 0 else str(x)
        df.loc[:, df.dtypes == float] = (
            df.loc[:, df.dtypes == float].applymap(format_int))

        # Replace {__sheet__} with sheet name
        df['__sheet__'] = sheet


        # Merge contestant database
        if database:
            alphanumeric = lambda text: re.sub(r'[^\w]', '', text).lower()
            process_str = lambda series: (
                series.apply(alphanumeric) if(series.dtype.kind == 'O')
                else series)

            for table in database.values():
                df_cols = df.columns.tolist()
                db_cols = table.columns.tolist()

                anchor_col = db_cols[0]
                overlapped_cols = [col for col in db_cols
                                   if (col in df_cols) and (col != anchor_col)]

                if anchor_col not in df_cols:  # TODO: Add a warning
                    continue

                # Copy processed values of anchor column to '__merge_anchor__'
                df['__merge_anchor__'] = process_str(df[anchor_col])
                table['__merge_anchor__'] = process_str(table[anchor_col])
                table = table.drop(columns=anchor_col)

                # Merge
                df = df.merge(table, on='__merge_anchor__', how='left')
                for col in overlapped_cols:
                    df[col] = df[f'{col}_y'].fillna(df[f'{col}_x'])
                    df = df.drop(columns=[f'{col}_x', f'{col}_y'])

                df = df.drop(columns='__merge_anchor__')


        if '__uid__' not in df.columns.tolist():
            avatar_mode = False

        # Fill in missing templates
        df['__template__'] = df['__template__'].fillna(1)
        df['__uid__'] = df['__uid__'].str.replace('_', '').str.strip()

        groups[sheet] = df

        if len(groups) == 1:
            print('\n\nHere is a snippet of your processed data:')

        print('\n' + preview_df(df, n_cols=len(df.columns), highlight=False))

    if not groups:
        Error(68).throw()


# Section G: Generate PowerPoint slides
    # Kill all PowerPoint instances
    run('TASKKILL /F /IM powerpnt.exe', stdout=DEVNULL, stderr=DEVNULL)


    OUTPUT_DIR = abs_path('output')
    AVATAR_DIR = abs_path('avatars')
    TEMP_DIR = abs_path('.temp')

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    os.makedirs(AVATAR_DIR, exist_ok=True)
    os.makedirs(TEMP_DIR, exist_ok=True)

    check_call(['attrib', '+H', TEMP_DIR])


    # Create Module1.bas
    if not os.path.exists(abs_path(TEMP_DIR, 'Module1.bas')):
        with open(abs_path(TEMP_DIR, 'Module1.bas'), 'w') as f:
            f.write(module1_bas)


    # Clear cache (TODO)
    with open(abs_path('.temp', 'last_clear_avatar_cache.txt'), 'w') as f:
        try:
            last_clear = int(f.readline())
        except:
            last_clear = 0

        if time.time() - last_clear > 1800:  # Clear every hour
            for avatar_path in os.scandir(AVATAR_DIR):
                os.unlink(avatar_path)

            # Update last clear time
            f.write(str(int(time.time())))


    # Download avatars with parallel processing
    if avatar_mode:
        print('\n\nDownloading avatars...')
        print('Make sure your internet connection is stable while we are '
              'downloading.')

        attempt = 0
        failed = False
        pool = Pool(min(4, len(token_list) + 1))
        uid_unknown_list = []

        while True:
            uid_list = []

            for df in groups.values():
                if df['__uid__'].dtype.kind in 'biufc':
                    Error(70).throw()

                for id in df['__uid__']:
                    if not (pd.isnull(id)
                            or os.path.isfile(get_avatar_path(AVATAR_DIR, id))
                            or id in uid_list
                            or id in uid_unknown_list):
                        uid_list.append(id)

            if not uid_list:
                break

            if attempt > 3:
                failed = True
                break

            try:
                pool.starmap(
                    download_avatar, 
                    zip(uid_list,
                        [AVATAR_DIR] * len(uid_list),
                        # Distribute the tokens among the user IDs
                        itertools.islice(
                            itertools.cycle(token_list), len(uid_list))
                    )
                )

            except (ConnectionError, TimeoutError) as e:
                if attempt == 3:
                    Error(20).throw(err_type=ErrorType.WARNING)

            except InvalidTokenError as e:
                Error(21.1).throw(*e.args)

            except UnknownUserError as e:
                uid_unknown_list.append(e.args[0])

            except DiscordAPIError as e:
                Error(22).throw(*e.args)

            attempt += 1

        if not failed:
            print('\033[A\033[2K\033[A\033[2K' + 'Avatar download complete!')

        if uid_unknown_list:
            Error(23).throw(str(uid_unknown_list), err_type=ErrorType.WARNING)

        pool.close()
        pool.join()

    time.sleep(0.2)
    print('\n\nGenerating slides...')
    print('Please do not click on any PowerPoint window that may '
          'show up in the process.\n')

    for sheet, df in groups.items():
        bar = ProgressBar(
            8, title=sheet, max_title_length=max(map(len, groups.keys())))

        # Open template presentation
        bar.set_description('Opening template.pptm')
        ppt = win32com.client.Dispatch('PowerPoint.Application')
        ppt.Presentations.Open(abs_path('template.pptm'))
        bar.add()

        # Import macros
        bar.set_description('Importing macros')

        try:
            ppt.VBE.ActiveVBProject.VBComponents.Import(
                abs_path(TEMP_DIR, 'Module1.bas'))
        except com_error as e:
            if e.hresult == -2147352567:  # type: ignore
                Error(41).throw()  # Trust access not yet enabled
            else:
                raise e

        bar.add()

        # Duplicate slides
        bar.set_description('Duplicating slides')
        slides_count = ppt.Run('Count')

        # Duplicate slides
        for template in df.loc[:, '__template__']:
            if as_type(int, template) not in range(1, slides_count + 1):
                df_showcase = df[
                    ['__template__']
                    + [col for col in df.columns if col != '__template__']]

                Error(71).throw(
                    f'{Style.BRIGHT}Template ID:{Style.NORMAL}     {template}'
                    f'\n'
                    f'{Style.BRIGHT}Error in sheet:{Style.NORMAL}  {sheet}',
                    f'Please also inspect merging sheets (signified with '
                    f'underscores at the beginning) if the problem could '
                    f'not directly be found in the mentioned sheet.',

                    preview_df(df_showcase,
                               n_cols=1,
                               words_to_highlight=[template]),
                )

            ppt.Run('Duplicate', template)

        bar.add()

        # Delete template slides when done
        ppt.Run('DelSlide', *range(1, slides_count + 1))
        bar.add()

        # Save as output file
        bar.set_description('Saving templates')
        output_path = abs_path(OUTPUT_DIR, f'{sheet}.pptx')

        ppt.Run('SaveAs', output_path)
        bar.add()

        run('TASKKILL /F /IM powerpnt.exe', stdout=DEVNULL, stderr=DEVNULL)
        bar.add()

        # Replace text
        bar.set_description('Filling in judging data')
        prs = Presentation(output_path)

        for i, slide in enumerate(prs.slides):
            replace_text(slide, df, i, avatar_mode)
        bar.add()

        # Save
        bar.set_description(f'Saving as {output_path}')
        prs.save(output_path)
        bar.add()


# Section H: Launch the file
    print(f'\nExported to {OUTPUT_DIR}')

    # Enable QuickEdit
    kernel32.SetConsoleMode(
        kernel32.GetStdHandle(-10), (0x4|0x80|0x20|0x2|0x10|0x1|0x40|0x100))

    inp('Press Enter to open the output folder...')
    os.startfile(OUTPUT_DIR)
